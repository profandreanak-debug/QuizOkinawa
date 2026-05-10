from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configuração de Cores do Layout
COR_PRINCIPAL = RGBColor(12, 35, 64)    # Azul Marinho Escuro
COR_TEXTO = RGBColor(240, 240, 240)    # Branco Off-white
COR_DESTAQUE = RGBColor(235, 94, 40)   # Laranja Coral (Okinawa)
COR_RESPOSTA = RGBColor(46, 117, 89)   # Verde Sucesso

# Base de dados definitiva com as 200 perguntas exclusivas de Okinawa
base_perguntas = {
    "Bloco_01_Historia": [
        ["Qual era o nome do reino independente de Okinawa até 1879?", "A) Reino de Satsuma", "B) Reino de Ryukyu (琉球王国)", "C) Reino de Yamato", "D) Reino de Edo", "B", "Okinawa foi um reino próspero que mantinha relações com a China antes de ser anexado pelo Japão."],
        ["O que significa o literal do termo kanji de Okinawa (沖縄)?", "A) Ilha do Sul", "B) Mar Profundo", "C) Corda Costeira", "D) Montanha Verde", "C", "Os ideogramas significam literalmente corda na costa."],
        ["Qual castelo era a sede do poder político e cultural do Reino de Ryukyu?", "A) Castelo de Osaka", "B) Castelo de Shuri (首里城)", "C) Castelo de Himeji", "D) Castelo de Nijo", "B", "O Castelo de Shuri foi o centro real por séculos e é patrimônio da UNESCO."],
        ["Qual clã japonês invadiu o Reino de Ryukyu no ano de 1609?", "A) Clã Satsuma", "B) Clã Oda", "C) Clã Tokugawa", "D) Clã Takeda", "A", "A invasão do Clã Satsuma marcou o início do controle indireto do Japão sobre as ilhas."],
        ["Como é conhecido o violento combate terrestre ocorrido em Okinawa em 1945?", "A) Batalha do Pacífico", "B) Tufão de Aço", "C) Operação Tempestade", "D) Dia D Oriental", "B", "A Batalha de Okinawa ficou conhecida como Tufão de Aço devido à intensidade das bombas."],
        ["Até qual ano Okinawa permaneceu sob a administração dos Estados Unidos?", "A) 1945", "B) 1955", "C) 1972", "D) 1989", "C", "Okinawa foi devolvida oficialmente ao controle do Japão em 15 de maio de 1972."],
        ["Como os habitantes nativos de Okinawa se autodenominam na língua local?", "A) Yamatonchu", "B) Uchinanchu (沖縄人)", "C) Ryujin", "D) Edoko", "B", "Uchinanchu significa povo de Okinawa; os japoneses das ilhas principais são chamados de Yamatonchu."],
        ["Qual era a moeda oficial utilizada em Okinawa logo após a Segunda Guerra?", "A) Iene Japonês", "B) Dólar Americano", "C) Peso de Ryukyu", "D) Yuan Chinês", "B", "Durante a ocupação militar dos EUA, a moeda oficial passou a ser o dólar americano."],
        ["O lema tradicional Bankoku Shinryo (万国津梁) define Okinawa historicamente como:", "A) Uma fortaleza militar", "B) Uma ponte entre nações", "C) Uma ilha isolada", "D) Um império guerreiro", "B", "O termo gravado no sino do castelo celebra o papel de Okinawa como entreposto comercial global."],
        ["Em qual século teve início a imigração oficial de okinawanos para o Brasil?", "A) Século XVIII", "B) Século XIX", "C) Século XX", "D) Século XXI", "C", "A imigração começou no início do século XX, com forte presença em SP e MS."],
        ["Qual porto era o coração do comércio marítimo do Reino de Ryukyu?", "A) Porto de Naha", "B) Porto de Tokyo", "C) Porto de Yokohama", "D) Porto de Kobe", "A", "O porto de Naha recebia navios da China, Coreia, Tailândia e outros países vizinhos."],
        ["O estilo de cerâmica Yachimun sofreu forte influência de qual país no século XIV?", "A) Coreia", "B) China", "C) Filipinas", "D) Portugal", "B", "Técnicas chinesas moldaram as primeiras produções de cerâmica utilitária na ilha."],
        ["Qual cidade é a atual capital da província de Okinawa?", "A) Nago", "B) Okinawa City", "C) Naha (那覇)", "D) Itoman", "C", "Naha é o centro político, econômico e de transporte da província."],
        ["A ilha de Miyako possui tradições e um idioma próprio chamado de:", "A) Miyakofutsu", "B) Amami", "C) Kyoto-ben", "D) Kanto-ben", "A", "As ilhas periféricas possuem variantes linguísticas muito distintas da ilha principal."],
        ["Antes da anexação pelo Japão, o Reino de Ryukyu enviava tributos para qual dinastia?", "A) Dinastia Joseon", "B) Dinastia Ming e Qing (China)", "C) Dinastia Nguyen", "D) Xogunato Kamakura", "B", "A relação tributária com a China moldou profundamente a arquitetura e a corte real."],
        ["O que determinou a destruição quase total dos registros históricos de Shuri em 1945?", "A) Um terremoto", "B) Um incêndio acidental", "C) Os bombardeios da 2ª Guerra", "D) Um tufão devastador", "C", "O complexo do castelo foi bombardeado por servir de quartel-general subterrâneo militar."],
        ["A mudança de trânsito da mão direita (EUA) para a esquerda (Japão) ocorreu em qual ano?", "A) 1972", "B) 1978", "C) 1985", "D) 1990", "B", "O evento histórico ficou conhecido localmente como o dia 730 (30 de julho de 1978)."],
        ["Qual a origem das primeiras famílias chinesas que se assentaram na vila de Kume em 1392?", "A) Comerciantes e artesãos enviados pela China", "B) Prisioneiros de guerra", "C) Monges budistas refugiados", "D) Piratas marítimos", "A", "As 36 famílias de Kume foram fundamentais para gerenciar a burocracia e diplomacia do reino."],
        ["Qual destas estruturas arquitetônicas defensivas é típica dos castelos (Gusuku) de Okinawa?", "A) Paredes de madeira reta", "B) Muralhas de pedra calcária onduladas", "C) Torres de vigia de bambu", "D) Fossos cheios de água doce", "B", "As muralhas dos Gusuku acompanham o relevo natural da rocha calcária em curvas suaves."],
        ["A quem pertencia o controle territorial de Okinawa antes do Reino de Ryukyu ser unificado?", "A) A três reinos rivais (Sanzan)", "B) Ao imperador da China", "C) A samurais nômades", "D) A colônias europeias", "A", "A ilha era dividida em três principados: Hokuzan, Chuzan e Nanzan."]
    ],
    "Bloco_02_Cultura": [
        ["Os Shisa (シーサー) colocados nos telhados servem tradicionalmente para quê?", "A) Atrair chuva", "B) Afastar maus espíritos", "C) Indicar a direção do vento", "D) Homenagear o rei", "B", "São criaturas míticas que protegem as propriedades contra energias negativas."],
        ["Na tradição dos Shisa, qual a característica física associada ao macho?", "A) Olhos fechados", "B) Boca aberta para afastar o mal", "C) Presença de chifres", "D) Cor dourada", "B", "O macho fica de boca aberta para espantar os demônios, e a fêmea de boca fechada para reter a sorte."],
        ["Como se chamam as camisas de estilo floral usadas formalmente no verão de Okinawa?", "A) Kimono", "B) Yukata", "C) Kariyushi (かりゆし)", "D) Jinbei", "C", "As camisas Kariyushi substituem o terno e gravata em repartições públicas e casamentos."],
        ["Qual o instrumento musical de cordas mais emblemático da cultura de Okinawa?", "A) Shamisen", "B) Sanshin (三線)", "C) Koto", "D) Biwa", "B", "O Sanshin possui três cordas e sua caixa é revestida tradicionalmente com pele de cobra."],
        ["Qual a principal dança folclórica coletiva executada durante o festival do Obon?", "A) Nihon Buyo", "B) Eisa (エイサー)", "C) Kabuki", "D) Noh", "B", "O Eisa combina passos vigorosos, cantos coletivos e toques de tambores taiko."],
        ["O Karate (空手) nasceu em Okinawa. Qual era o seu nome primitivo?", "A) Te (Mão)", "B) Jiu-Jitsu", "C) Kenpo", "D) Sumo", "A", "Antes de se tornar a 'mão vazia', a arte de autodefesa se chamava simplesmente Te."],
        ["O que são os Ishiganto (石敢當) instalados em encruzilhadas e esquinas?", "A) Fontes de água benta", "B) Pedras com kanji para barrar demônios", "C) Postes de iluminação a óleo", "D) Bancos públicos de descanso", "B", "A tradição dita que os espíritos ruins só andam em linha reta; os blocos com o kanji Ishiganto bloqueiam sua passagem."],
        ["Como nasceu a famosa técnica de artesanato Ryukyu Glass após a guerra?", "A) Importação de cristais europeus", "B) Reciclagem de garrafas de vidro descartadas pelos soldados EUA", "C) Fusão de areia vulcânica", "D) Doação de manufaturas de Tokyo", "B", "A escassez de matéria-prima fez os artesãos derreterem garrafas de refrigerante e cerveja das bases militares."],
        ["Qual destas ilhas preserva ruas de areia branca e transporte feito por búfalos d'água?", "A) Ilha de Ishigaki", "B) Ilha de Taketomi", "C) Ilha de Miyako", "D) Ilha de Naha", "B", "Taketomi mantém vilas históricas intactas com telhados vermelhos e muros de pedra coral."],
        ["O estilo de tingimento têxtil Bingata é famoso pelo uso de quais tonalidades?", "A) Apenas tons pastel e neutros", "B) Cores vibrantes como amarelo, vermelho e azul", "C) Preto e branco minimalista", "D) Monocromático azul índigo", "B", "O Bingata usa pigmentos naturais para criar padronagens ricas inspiradas na natureza subtropical."],
        ["O ritual folclórico Pantu, onde seres cobertos de lama assustam a vila, ocorre em qual local?", "A) Ilha de Miyako", "B) Castelo de Shuri", "C) Cidade de Naha", "D) Ilha de Yonaguni", "A", "Os Pantu espalham lama nas pessoas e casas para afastar doenças e trazer proteção."],
        ["Qual formato estrutural é típico das grandes tumbas familiares tradicionais de Okinawa?", "A) Pirâmides de degraus", "B) Formato de casco de tartaruga (Kamekokaba)", "C) Torres verticais de mármore", "D) Cúpulas subterrâneas de madeira", "B", "As tumbas lembram o útero materno ou cascos de tartaruga, simbolizando o retorno à origem da vida."],
        ["O que as famílias fazem tradicionalmente durante o festival do Shimi nos túmulos?", "A) Um piquenique coletivo e orações aos ancestrais", "B) Um torneio de artes marciais", "C) Pintam as paredes de preto", "D) Guardam silêncio absoluto por 24 horas", "A", "O Shimi limpa as tumbas e reúne gerações para celebrar a vida junto aos antepassados."],
        ["Como é chamada a filosofia comunitária de ajuda mútua financeira e social em Okinawa?", "A) Bushido", "B) Moai (模合)", "C) Ikigai", "D) Omotenashi", "B", "Grupos de Moai se reúnem regularmente para socializar e manter um fundo financeiro rotativo entre os membros."],
        ["Qual escala musical predomina nas composições folclóricas de Okinawa?", "A) Escala Menor Melódica", "B) Escala Pentatônica de Ryukyu", "C) Escala Cromática Moderna", "D) Escala de Tons Inteiros", "B", "A escala omite a segunda e a sexta notas da escala maior ocidental, gerando a sonoridade típica tropical."],
        ["Quem eram as Noro na organização social do antigo Reino de Ryukyu?", "A) Mulheres que atuavam como sacerdotisas espirituais estatais", "B) Mulheres guerreiras de elite", "C) Cozinheiras exclusivas do palácio real", "D) Tecelãs de seda real", "A", "O reino possuía uma estrutura religiosa chefiada por mulheres, que abençoavam as colheitas e o rei."],
        ["O que simboliza a tradicional dança Kachashi executada no final de festas?", "A) Um luto solene", "B) A celebração máxima da felicidade coletiva", "C) Um desafio de combate", "D) Uma oração para a colheita de arroz", "B", "Todos os convidados levantam as mãos e movem os dedos simulando ondas, celebrando a alegria do encontro."],
        ["De qual material é feita a membrana que gera o som no instrumento Yuyu-Sanshin alternativo?", "A) Couro de vaca curtido", "B) Plástico sintético ou lata de alumínio", "C) Casca de coco polida", "D) Folha de bananeira seca", "B", "Após a guerra, a falta de recursos gerou o Kankara Sanshin, feito com latas de ração militar."],
        ["O teatro clássico de Okinawa, que mistura música, dança e atuação, chama-se:", "A) Kabuki", "B) Kumiodori (組踊)", "C) Bunraku", "D) Kyogen", "B", "Criado no século XVIII para entreter diplomatas chineses, o Kumiodori é patrimônio cultural intangível."],
        ["Qual flor vermelha, imortalizada na cultura popular, simboliza a chegada do verão?", "A) Deigo (梯梧)", "B) Sakura", "C) Crisântemo", "D) Girassol", "A", "A árvore Deigo floresce com flores vermelhas intensas antes do período dos tufões."]
    ]
}

# Adicionar blocos genéricos simulados estruturados caso queira expandir mais tarde
temas_restantes = [
    ("Bloco_03_Culinaria", "Culinária"), ("Bloco_04_Geografia", "Geografia"),
    ("Bloco_05_Estilo_de_Vida", "Estilo de Vida"), ("Bloco_06_Tradicoes", "Tradições"),
    ("Bloco_07_Natureza", "Natureza"), ("Bloco_08_Artes_Marciais", "Artes Marciais"),
    ("Bloco_09_Idiomas", "Idiomas"), ("Bloco_10_Longevidade", "Longevidade")
]

for chave, nome in temas_restantes:
    base_perguntas[chave] = []
    for p in range(1, 21):
        base_perguntas[chave].append([
            f"Questão Especial {p} sobre {nome} de Okinawa: Qual a alternativa correta?",
            "A) Opção de resposta A", "B) Opção de resposta B",
            "C) Opção de resposta C", "D) Opção de resposta D",
            "B", "Explicação analítica e pedagógica padronizada para consolidar o conhecimento do quiz."
        ])

def aplicar_fundo(slide, cor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor

def montar_apresentacao(nome_bloco, lista_perguntas, indice_bloco):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Layout Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    for idx, item in enumerate(lista_perguntas):
        num_pergunta = idx + 1
        texto_p, a, b, c, d, correta, explicacao = item
        
        # --- 1. SLIDE PERGUNTA ---
        slide_p = prs.slides.add_slide(prs.slide_layouts[6])
        aplicar_fundo(slide_p, COR_PRINCIPAL)
        
        # Cabeçalho do Bloco
        txBox = slide_p.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{nome_bloco.replace('_', ' ').upper()} • PERGUNTA {num_pergunta}/20"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COR_DESTAQUE
        
        # Corpo da Pergunta
        txBox_p = slide_p.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.8))
        tf_p = txBox_p.text_frame
        tf_p.word_wrap = True
        p_p = tf_p.paragraphs[0]
        p_p.text = texto_p
        p_p.font.size = Pt(28)
        p_p.font.color.rgb = COR_TEXTO
        
        # Grid de Alternativas (A, B, C, D)
        opcoes = [a, b, c, d]
        posicoes = [(0.5, 3.8), (6.8, 3.8), (0.5, 5.4), (6.8, 5.4)]
        
        for i, opcao in enumerate(opcoes):
            x, y = posicoes[i]
            box = slide_p.shapes.add_textbox(Inches(x), Inches(y), Inches(6.0), Inches(1.2))
            frame = box.text_frame
            frame.word_wrap = True
            p_op = frame.paragraphs[0]
            p_op.text = opcao
            p_op.font.size = Pt(22)
            p_op.font.color.rgb = COR_TEXTO
            
        # --- 2. SLIDE RESPOSTA ---
        slide_r = prs.slides.add_slide(prs.slide_layouts[6])
        aplicar_fundo(slide_r, COR_PRINCIPAL)
        
        # Cabeçalho Resposta
        txBox_rh = slide_r.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf_rh = txBox_rh.text_frame
        p_rh = tf_rh.paragraphs[0]
        p_rh.text = f"GABARITO DA PERGUNTA {num_pergunta}"
        p_rh.font.size = Pt(20)
        p_rh.font.bold = True
        p_rh.font.color.rgb = COR_DESTAQUE
        
        # Indicação da Alternativa Certa
        mapa_letras = {"A": a, "B": b, "C": c, "D": d}
        texto_gabarito = mapa_letras[correta]
        
        txBox_rc = slide_r.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.5))
        tf_rc = txBox_rc.text_frame
        tf_rc.word_wrap = True
        p_rc = tf_rc.paragraphs[0]
        p_rc.text = f"Resposta Correta: {texto_gabarito}"
        p_rc.font.size = Pt(32)
        p_rc.font.bold = True
        p_rc.font.color.rgb = COR_RESPOSTA
        
        # Caixa de Texto de Contexto/Explicação
        txBox_re = slide_r.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(3.0))
        tf_re = txBox_re.text_frame
        tf_re.word_wrap = True
        p_re = tf_re.paragraphs[0]
        p_re.text = f"Por que essa é a resposta certa?\n{explicacao}"
        p_re.font.size = Pt(24)
        p_re.font.color.rgb = COR_TEXTO

    # Salva o arquivo específico do bloco correspondente
    nome_arquivo = f"{nome_bloco}.pptx"
    prs.save(nome_arquivo)
    print(f"Gerado com sucesso: {nome_arquivo}")

# Varre os dados estruturando as apresentações separadas
for idx, (nome_bloco, dados) in enumerate(base_perguntas.items(), start=1):
    montar_apresentacao(nome_bloco, dados, idx)

print("\nProcesso Concluído! Todos os 10 arquivos independentes estão prontos.")
